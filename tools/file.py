import os
import pandas as pd
from docx import Document
from pptx import Presentation
from PIL import Image
from tools.base import BaseTool

class FileTool(BaseTool):
    name = "file"
    description = "读取或写入文件。支持格式：txt, md, json, csv, xlsx, docx, pptx, jpg, png (read-only info)。操作：read, write。"
    args_schema = {
        "operation": "string (read | write)",
        "path": "string",
        "content": "string (optional, for write)"
    }

    def run(self, operation: str, path: str, content: str = None) -> str:
        operation = operation.lower()
        path = path.strip()
        
        if operation == "read":
            return self._read_file(path)
        elif operation == "write":
            if content is None:
                return "Error: 'content' is required for write operation."
            return self._write_file(path, content)
        else:
            return f"Error: Unknown operation '{operation}'. Use 'read' or 'write'."

    def _read_file(self, path: str) -> str:
        if not os.path.exists(path):
            return f"Error: File '{path}' not found."
            
        ext = os.path.splitext(path)[1].lower()
        
        try:
            if ext in ['.xlsx', '.xls']:
                # Read Excel
                df = pd.read_excel(path)
                return self._dataframe_to_string(df, path)
            
            elif ext in ['.csv']:
                # Read CSV
                df = pd.read_csv(path)
                return self._dataframe_to_string(df, path)
                
            elif ext in ['.docx']:
                # Read Word
                doc = Document(path)
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                return '\n'.join(full_text)
                
            elif ext in ['.pptx']:
                # Read PowerPoint
                prs = Presentation(path)
                text_content = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)
                return '\n'.join(text_content)
                
            elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif']:
                # Read Image Info
                with Image.open(path) as img:
                    info = f"Image File: {path}\nFormat: {img.format}\nSize: {img.size}\nMode: {img.mode}"
                    return info
            
            else:
                # Default to text read
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
                    
        except Exception as e:
            return f"Error reading file '{path}': {str(e)}"

    def _dataframe_to_string(self, df, path):
        """Helper to format dataframe output with limits"""
        rows, cols = df.shape
        info = f"File: {path}\nShape: {rows} rows, {cols} columns\nColumns: {list(df.columns)}\n"
        
        if rows > 50:
            info += f"\n[Warning] File is large. Showing first 50 rows only to avoid context overflow.\n"
            return info + df.head(50).to_string()
        else:
            return info + df.to_string()

    def _write_file(self, path: str, content: str) -> str:
        ext = os.path.splitext(path)[1].lower()
        
        try:
            # Create directory if needed
            os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
            
            if ext in ['.xlsx', '.xls']:
                import io
                import json
                
                df = None
                try:
                    data = json.loads(content)
                    df = pd.DataFrame(data)
                except:
                    try:
                        df = pd.read_csv(io.StringIO(content))
                    except:
                        pass
                
                if df is not None:
                    df.to_excel(path, index=False)
                    return f"Successfully wrote Excel file to {path}"
                else:
                    return "Error: Content for Excel must be valid JSON list-of-dicts or CSV string."

            elif ext in ['.docx']:
                doc = Document()
                doc.add_paragraph(content)
                doc.save(path)
                return f"Successfully wrote Word file to {path}"
                
            else:
                with open(path, 'w', encoding='utf-8') as f:
                    f.write(content)
                return f"Successfully wrote file to {path}"
                
        except Exception as e:
            return f"Error writing file '{path}': {str(e)}"
